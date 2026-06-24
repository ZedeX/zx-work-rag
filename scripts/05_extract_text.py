"""Step 5: Extract text from all text-category files.

Supports: docx, pptx, xlsx, pdf, txt, csv, html, htm, xml, json
For converted files, uses the converted_path.
For non-text files (image/video/audio), creates a metadata-only record.
"""
import sys
import sqlite3
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from config.settings import config, Config
from loguru import logger


def setup_logging():
    log_file = Path(config.LOG_DIR) / "05_extract_text.log"
    log_file.parent.mkdir(parents=True, exist_ok=True)
    logger.add(str(log_file), rotation="10 MB", encoding="utf-8",
               format="{time:YYYY-MM-DD HH:mm:ss} | {level} | {message}")


def extract_docx(filepath: str) -> str:
    """Extract text from .docx files."""
    from docx import Document
    doc = Document(filepath)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                parts.append(" | ".join(row_text))
    return "\n".join(parts)


def extract_pptx(filepath: str) -> str:
    """Extract text from .pptx files."""
    from pptx import Presentation
    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_parts.append(para.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_parts.append(" | ".join(row_text))
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_xlsx(filepath: str) -> str:
    """Extract text from .xlsx files."""
    from openpyxl import load_workbook
    wb = load_workbook(filepath, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts = [f"--- Sheet: {sheet_name} ---"]
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if row_text:
                sheet_parts.append(" | ".join(row_text))
                row_count += 1
            if row_count >= 500:  # Limit rows per sheet
                sheet_parts.append(f"... (truncated at 500 rows, total {ws.max_row})")
                break
        parts.append("\n".join(sheet_parts))
    wb.close()
    return "\n\n".join(parts)


def extract_pdf(filepath: str) -> str:
    """Extract text from .pdf files using pdfplumber."""
    import pdfplumber
    parts = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                parts.append(text.strip())
            # Also extract tables
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    row_text = [str(cell).strip() for cell in row if cell and str(cell).strip()]
                    if row_text:
                        parts.append(" | ".join(row_text))
    return "\n".join(parts)


def extract_txt(filepath: str) -> str:
    """Extract text from plain text files."""
    encodings = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]
    for enc in encodings:
        try:
            with open(filepath, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return ""


def extract_html(filepath: str) -> str:
    """Extract text from HTML files."""
    from bs4 import BeautifulSoup
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    # Remove script and style elements
    for s in soup(["script", "style"]):
        s.decompose()
    return soup.get_text(separator="\n", strip=True)


def extract_csv(filepath: str) -> str:
    """Extract text from CSV files."""
    return extract_txt(filepath)


def extract_json(filepath: str) -> str:
    """Extract text from JSON files."""
    import json
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, ensure_ascii=False, indent=2)[:50000]
    except Exception:
        return extract_txt(filepath)


# Dispatch table: extension -> extractor function
EXTRACTORS = {
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".pdf": extract_pdf,
    ".txt": extract_txt,
    ".csv": extract_csv,
    ".html": extract_html,
    ".htm": extract_html,
    ".xml": extract_txt,
    ".json": extract_json,
    ".md": extract_txt,
    ".log": extract_txt,
    ".rtf": extract_txt,
}


def extract_text(filepath: str, ext: str) -> tuple:
    """Extract text from a file. Returns (text, success, error)."""
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return "", False, f"No extractor for {ext}"

    try:
        text = extractor(filepath)
        if text and text.strip():
            return text, True, ""
        else:
            return "", False, "Empty text extracted"
    except Exception as e:
        return "", False, str(e)[:200]


def create_metadata_record(filepath: str, folder_path: str, folder_tags: str,
                           file_name: str, ext: str, size: int) -> str:
    """Create a metadata-only text record for non-text files (images, videos, audio).

    This allows these files to be found by folder/tag search even though
    their content isn't extracted.
    """
    parts = [
        f"[File Metadata Record]",
        f"File name: {file_name}",
        f"File type: {ext or 'unknown'}",
        f"File size: {size} bytes",
        f"Location: {folder_path}",
    ]
    if folder_tags:
        parts.append(f"Tags: {folder_tags}")
    return "\n".join(parts)


def add_text_columns(conn: sqlite3.Connection):
    """Add text-related columns if not exist."""
    c = conn.cursor()
    for col, col_type in [
        ("extracted_text", "TEXT"),
        ("text_length", "INTEGER DEFAULT 0"),
        ("converted_path", "TEXT DEFAULT ''"),
    ]:
        try:
            c.execute(f"ALTER TABLE files ADD COLUMN {col} {col_type}")
        except sqlite3.OperationalError:
            pass
    conn.commit()


def extract_all(conn: sqlite3.Connection):
    """Extract text from all eligible files."""
    c = conn.cursor()

    # Get text-category files that haven't been extracted yet
    c.execute("""
        SELECT id, file_path, extension, file_category, folder_path,
               folder_tags, file_name, size, converted_path, status
        FROM files
        WHERE is_duplicate = 0
        AND text_extracted = 0
        AND status IN ('scanned', 'identified', 'converted')
        ORDER BY file_category, extension
    """)
    rows = c.fetchall()
    logger.info(f"Files to process: {len(rows)}")

    text_extracted = 0
    metadata_only = 0
    failed = 0

    for row in rows:
        (file_id, filepath, ext, category, folder_path,
         folder_tags, file_name, size, converted_path, status) = row

        # For non-text files, create metadata-only record
        if category in ("image", "video", "audio", "archive", "code", "other"):
            meta_text = create_metadata_record(
                filepath, folder_path, folder_tags, file_name, ext, size
            )
            c.execute("""
                UPDATE files
                SET extracted_text = ?, text_length = ?, text_extracted = 1, status = 'metadata_only'
                WHERE id = ?
            """, (meta_text, len(meta_text), file_id))
            metadata_only += 1
            continue

        # For converted files, use the converted path
        actual_path = converted_path if converted_path and Path(converted_path).exists() else filepath
        actual_ext = ext

        # If file was converted, use the new extension
        if converted_path and Path(converted_path).exists():
            actual_ext = Path(converted_path).suffix.lower()

        if not Path(actual_path).exists():
            c.execute("UPDATE files SET status = 'missing' WHERE id = ?", (file_id,))
            continue

        text, ok, err = extract_text(actual_path, actual_ext)

        if ok:
            # Truncate very long texts to 100KB
            if len(text) > 100000:
                text = text[:100000] + "\n... (truncated)"
            c.execute("""
                UPDATE files
                SET extracted_text = ?, text_length = ?, text_extracted = 1, status = 'text_extracted'
                WHERE id = ?
            """, (text, len(text), file_id))
            text_extracted += 1
        else:
            c.execute("""
                UPDATE files
                SET text_extracted = 0, status = 'extraction_failed', error_message = ?
                WHERE id = ?
            """, (err, file_id))
            failed += 1
            if failed <= 50:
                logger.warning(f"Extraction failed: {filepath} - {err}")

        total = text_extracted + metadata_only + failed
        if total % 500 == 0:
            conn.commit()
            logger.info(f"Text: {text_extracted}, Meta: {metadata_only}, Failed: {failed} / {len(rows)}")

    conn.commit()
    logger.info(f"Extraction complete. Text: {text_extracted}, Meta-only: {metadata_only}, Failed: {failed}")

    # Summary
    c.execute("SELECT status, COUNT(*) FROM files WHERE is_duplicate = 0 GROUP BY status ORDER BY COUNT(*) DESC")
    logger.info("=== Status Summary ===")
    for row in c.fetchall():
        logger.info(f"  {row[0]}: {row[1]}")


def main():
    setup_logging()
    logger.info("=" * 60)
    logger.info("Step 5: Text Extraction")
    logger.info("=" * 60)

    conn = sqlite3.connect(config.DB_PATH)
    add_text_columns(conn)
    extract_all(conn)
    conn.close()
    logger.info("Done.")


if __name__ == "__main__":
    main()
